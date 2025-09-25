#App.py
import streamlit as st
import pandas as pd
import json
import matplotlib.pyplot as plt
import plotly.express as px
import streamlit.components.v1 as components
from langchain_core.messages import HumanMessage, AIMessage
import os
import base64
from fuzzysearch import find_near_matches
import requests
from dotenv import load_dotenv
from PIL import Image
import docx2txt
import tempfile
from pptx import Presentation
from multi_agent import (
    metadata_agent_executor,
    analytics_executor,
    data_entry_executor,
    event_data_executor,
    tracker_data_executor,
    # azure_document_intelligence_executor,
    routing_decision
)
# --- Import agent executors ---
# from agents.metadata_agent import metadata_agent_executor
# from agents.analytics_agent import analytics_executor
# from agents.tracker_data_agent import tracker_data_executor
# from agents.event_data_agent import event_data_executor
# from agents.data_entry_agent import data_entry_executor

from agents.tools.faiss_search.embed_dhis2_faiss_metadata import run_embedding
from agents.tools.analytics_tools import get_all

load_dotenv()

# Load environment variables

DHIS2_BASE_URL = os.getenv("DHIS2_BASE_URL")
DHIS2_USERNAME = os.getenv("DHIS2_USERNAME")
DHIS2_PASSWORD = os.getenv("DHIS2_PASSWORD")


# ========== Helpers ==========
def generate_distinct_colors(n):
    return [f"hsl({i * 360 / n}, 70%, 50%)" for i in range(n)]


def get_data(url, doc_type, disaggregations, indicators, periods, org_units):
    flattened = []
    org_units_lookup = []
    indicator_string = ";".join(indicators)
    period_string = ";".join(periods)
    org_unit_string = ";".join(org_units)

    if len(org_units) > 0:
        get_the_org = get_all(
            endpoint="organisationUnits.json",
            key="organisationUnits",
            fields="id,name",
            # fields="id,code,name,parent,children,level,path,ancestors",
            filters={"id": f"in:[{','.join(org_units)}]"}
        )

        for org in get_the_org:
            org_units_lookup.append({
                "org": org["name"],
                "org_id": org["id"],
            })

    dimensions = [
        f"dx:{indicator_string}",
        f"pe:{period_string}",
        f"ou:{org_unit_string}"
    ]
    if doc_type == "indicator":
        include_num_den = True
        include_coc_dimension = False
    elif doc_type == "dataElement":
        get_the_de = get_all(
            endpoint="dataElements.json",
            key="dataElements",
            fields="id,name,categoryCombo[id,name,categories[id,name,categoryOptions[id,name]]]",
            # fields="id,code,name,parent,children,level,path,ancestors",
            filters={"id": f"in:[{','.join(indicators)}]"}
        )

        for de in get_the_de:
            categories = de.get("categoryCombo", {}).get("categories", {})
            # Flatten the structure
            for category in categories:
                for option_ in category["categoryOptions"]:
                    flattened.append({
                        "dataElement": de["name"],
                        "dataElement_id": de["id"],
                        "category": category["name"],
                        "category_id": category["id"],
                        "option": option_["name"],
                        "option_id": option_["id"]
                    })

        # After flattening is done
        unique_categories = sorted({item["category"] for item in flattened})
        print("✅ Unique categories found:")
        for cat in unique_categories:
            print(f"- {cat}")

        # Filter flattened items whose category name matches any disaggregation
        if disaggregations is not None:
            filtered_flattened = []
            for item in flattened:
                category_name = item["category"].lower().strip()
                for disaggregation in disaggregations:
                    matches = find_near_matches(disaggregation.lower().strip(), category_name, max_l_dist=2)
                    if matches:
                        filtered_flattened.append(item)
                        break  # Stop after the first match for this item
            # Build dictionary of {category_id: [option_id, ...]}
            disaggregations_dict = {}
            for item in filtered_flattened:
                cat_id = item["category_id"]
                opt_id = item["option_id"]
                disaggregations_dict.setdefault(cat_id, []).append(opt_id)

            # Construct dimensions list
            dimensions = [
                f"dx:{indicator_string}",
                f"pe:{period_string}",
                f"ou:{org_unit_string}"
            ]

            # Add disaggregation dimensions
            for cat_id, option_ids in disaggregations_dict.items():
                option_string = ";".join(sorted(set(option_ids)))  # Optional: deduplicate + sort
                dimensions.append(f"{cat_id}:{option_string}")

            dimensions.append("co")  # ⚠️ Only add this if dx supports category option combos

        include_num_den = False
        include_coc_dimension = True  # Often useful for DEs

        # Add co dimension explicitly if requested
        # if include_coc_dimension:
        #     dimensions.append("co")  # ⚠️ Only add this if dx supports category option combos


    elif doc_type == "programIndicator":
        include_num_den = False
        include_coc_dimension = True  # Often useful for DEs
        pass
    else:
        include_num_den = False
        include_coc_dimension = True  # Often useful for DEs
        pass
    params = {
        "dimension": dimensions,
        # "displayProperty": display_property,
        "includeNumDen": str(include_num_den).lower(),
        "skipMeta": "false",
        "skipData": "false"
    }
    url = f"{DHIS2_BASE_URL.rstrip('/')}/api/analytics"
    try:
        response = requests.get(
            url=url,
            auth=(DHIS2_USERNAME, DHIS2_PASSWORD),
            params=params
        )
        response.raise_for_status()
        return {
            "meta_data": response.json(),
            "org_units_lookup": org_units_lookup,
            "flattened": flattened,
        }
    except Exception as e:
        return {"error": str(e)}


def analytics_to_dataframe(tool_data: dict):
    if "data" not in tool_data:
        return pd.DataFrame()
    url = tool_data["url"]
    doc_type = tool_data["doc_type"]
    disaggregations = tool_data["disaggregations"]
    indicators = tool_data["indicators"]
    periods = tool_data["periods"]
    org_units = tool_data["org_units"]
    url = url.replace("skipMeta=true", "skipMeta=false")
    meta_data = get_data(url, doc_type, disaggregations, indicators, periods, org_units)
    if isinstance(meta_data, str):
        meta_data = json.loads(meta_data)

    # Now access it safely
    data = meta_data.get("meta_data", {})
    org_units_lookup = meta_data.get("org_units_lookup", {}),
    flattened = meta_data.get("flattened", {}),

    headers = data.get("headers", [])
    rows = data.get("rows", [])
    columns = [header["name"] for header in headers]
    return org_units_lookup, flattened, pd.DataFrame(rows, columns=columns)


def render_chart_chartjs(filtered_df: pd.DataFrame, selected_indicators: list, chart_type: str):
    if filtered_df.empty or not selected_indicators:
        st.info("No data available or indicators selected.")
        return

    this_df = filtered_df.rename(columns={
        "Data": "dx",
        "Period": "pe",
        "Organisation unit": "ou",
        "Value": "value"
    })
    # # Handle optional category option combo column
    # if "co" in df.columns:
    #     df = df.rename(columns={"co": "Category option combo"})

    this_df["value"] = pd.to_numeric(this_df["value"], errors="coerce")

    grouped = this_df[this_df["dx"].isin(selected_indicators)]
    # Define available label columns
    available_label_columns = ["pe", "ou"] + [col for col in this_df.columns if col.startswith("co_")]

    # Ensure "pe" is the default (index=0)
    label_column = st.selectbox("Choose x-axis (label) dimension:", options=available_label_columns,
                                index=available_label_columns.index("pe"))

    labels = sorted(grouped[label_column].dropna().unique().tolist())
    colors = ["#4285F4", "#DB4437", "#F4B400", "#0F9D58", "#AB47BC"]
    datasets = []

    for idx, indicator in enumerate(selected_indicators):
        subset = grouped[grouped["dx"] == indicator]
        subset_grouped = subset.groupby(label_column)["value"].sum().reindex(labels).fillna(0)
        values = subset_grouped.tolist()
        dynamic_colors = generate_distinct_colors(len(values))
        if chart_type == "pie":
            datasets.append({
                "label": indicator,
                "data": values,
                "backgroundColor": dynamic_colors,
                "borderColor": "#fff",
                "borderWidth": 1
            })
        else:
            datasets.append({
                "label": indicator,
                "data": values,
                "backgroundColor": dynamic_colors,
                "borderColor": dynamic_colors,
                "fill": chart_type != "line",
                "tension": 0.3
            })

    chart_config = {
        "type": chart_type,
        "data": {
            "labels": labels,
            "datasets": datasets
        },
        "options": {
            "responsive": True,
            "maintainAspectRatio": False,
            "plugins": {
                "legend": {"position": "top"},
                "title": {"display": True, "text": "DHIS2 Analytics Chart"},
                "datalabels": {
                    "anchor": "end",
                    "align": "top",
                    "color": "#000",
                    "font": {"weight": "bold"},
                    "formatter": "Math.round"
                }
            },
            "scales": {
                "y": {
                    "beginAtZero": True
                }
            }
        },
        "plugins": ["ChartDataLabels"]
    }

    chart_json = json.dumps(chart_config)

    # Encode CSV data as base64 for safe injection
    csv_bytes = this_df.to_csv(index=False).encode("utf-8")
    csv_b64 = base64.b64encode(csv_bytes).decode("utf-8")

    html_code = f"""
    <script src="https://cdn.jsdelivr.net/npm/chart.js"></script>
    <script src="https://cdn.jsdelivr.net/npm/chartjs-plugin-datalabels@2"></script>

    <div style="position: relative; height: 400px;">
        <canvas id="myChart"></canvas>
    </div>

    <div style="margin-top: 10px;">
        <button onclick="downloadChart()" style="padding: 8px 16px; font-size: 14px; margin-right: 10px;">
            📥 Download Chart as PNG
        </button>
        <button onclick="downloadCSV()" style="padding: 8px 16px; font-size: 14px;">
            📥 Download Data as CSV
        </button>
    </div>

    <script>
    const ctx = document.getElementById('myChart').getContext('2d');
    Chart.register(ChartDataLabels);
    new Chart(ctx, {chart_json});

    function downloadChart() {{
        const link = document.createElement('a');
        link.download = 'chart.png';
        link.href = document.getElementById('myChart').toDataURL('image/png');
        link.click();
    }}

    function downloadCSV() {{
        const csvData = atob("{csv_b64}");
        const blob = new Blob([csvData], {{ type: 'text/csv;charset=utf-8;' }});
        const url = URL.createObjectURL(blob);
        const link = document.createElement('a');
        link.setAttribute('href', url);
        link.setAttribute('download', 'filtered_data.csv');
        document.body.appendChild(link);
        link.click();
        document.body.removeChild(link);
    }}
    </script>
    """

    st.subheader("📊 Chart Preview")
    components.html(html_code, height=500)


def render_chart_matplotlib(df: pd.DataFrame, selected_indicators: list, chart_type: str):
    if df.empty or not selected_indicators:
        st.info("No data available or indicators selected.")
        return

    fig, ax = plt.subplots(figsize=(10, 5))
    grouped = df[df["dx"].isin(selected_indicators)]

    if chart_type == "bar":
        for indicator in selected_indicators:
            subset = grouped[grouped["dx"] == indicator]
            ax.bar(subset["pe"], subset["value"], label=indicator)
    elif chart_type == "line":
        for indicator in selected_indicators:
            subset = grouped[grouped["dx"] == indicator]
            ax.plot(subset["pe"], subset["value"], label=indicator, marker='o')
    elif chart_type == "pie" and len(selected_indicators) == 1:
        subset = grouped[grouped["dx"] == selected_indicators[0]]
        subset = subset.groupby("pe")["value"].sum()
        ax.pie(subset.values, labels=subset.index, autopct='%1.1f%%')
    else:
        st.warning("Pie chart only supports one indicator.")
        return

    if chart_type != "pie":
        ax.set_xlabel("Period")
        ax.set_ylabel("Value")
        ax.set_title("DHIS2 Analytics Chart (Matplotlib)")
        ax.legend()

    st.subheader("📊 Chart Preview (Matplotlib)")
    st.pyplot(fig)


def render_chart_plotly(df: pd.DataFrame, selected_indicators: list, chart_type: str):
    if df.empty or not selected_indicators:
        st.info("No data available or indicators selected.")
        return

    filtered_df = df[df["dx"].isin(selected_indicators)]

    if chart_type == "bar":
        fig = px.bar(filtered_df, x="pe", y="value", color="dx", barmode="group",
                     title="DHIS2 Analytics Chart (Plotly)")
    elif chart_type == "line":
        fig = px.line(filtered_df, x="pe", y="value", color="dx", markers=True, title="DHIS2 Analytics Chart (Plotly)")
    elif chart_type == "pie" and len(selected_indicators) == 1:
        subset = filtered_df[filtered_df["dx"] == selected_indicators[0]]
        pie_data = subset.groupby("pe")["value"].sum().reset_index()
        fig = px.pie(pie_data, names="pe", values="value", title=f"{selected_indicators[0]} (Pie Chart)")
    else:
        st.warning("Pie chart only supports one indicator.")
        return

    st.subheader("📊 Chart Preview (Plotly)")
    st.plotly_chart(fig, use_container_width=True)


def chart_data(result_, chart_backend):
    raw_data = result_.get("raw_data", {})

    if raw_data is None:
        st.warning("No visualization data returned from the analytics executor.")
        return
    org_units_lookup, flattened, df_chart_data = analytics_to_dataframe(raw_data)

    st.session_state.raw_data_df = df_chart_data

    raw_data_df = df_chart_data

    # Base required columns
    base_required_cols = {"dx", "pe", "value"}
    # Columns to exclude (known standard fields)
    excluded_cols = {"dx", "pe", "ou", "value", "numerator", "denominator", "factor", "multiplier", "divisor", "co"}
    # Detect up to 3 disaggregation columns (e.g., category option combos)
    disagg_cols = [col for col in raw_data_df.columns if col not in excluded_cols][:3]

    # Add disaggregation columns to required cols
    required_cols = base_required_cols.union(disagg_cols)

    if not raw_data_df.empty and required_cols.issubset(raw_data_df.columns):
        with st.container():
            # If org_units_lookup or flattened are strings, convert them to Python objects
            if isinstance(org_units_lookup, str):
                org_units_lookup = json.loads(org_units_lookup)

            if isinstance(flattened, str):
                flattened = json.loads(flattened)

            # If they are a tuple (e.g., result of a return statement like `return x,`), unwrap them
            if isinstance(org_units_lookup, tuple):
                org_units_lookup = org_units_lookup[0]

            if isinstance(flattened, tuple):
                flattened = flattened[0]

            org_map = {item['org_id']: item['org'] for item in org_units_lookup}
            dx_map = {item['dataElement_id']: item['dataElement'] for item in flattened}
            option_map = {item['option_id']: item['option'] for item in flattened}
            category_map = {item['category_id']: item['category'] for item in flattened}

            # --- Step 2: Replace IDs with names in DataFrame (only if columns exist) ---
            def replace_if_exists(df: pd.DataFrame, column: str, mapping: dict):
                if column in df.columns:
                    df[column] = df[column].map(mapping).fillna(df[column])

            # Assume filtered_df is already defined
            # Replace `ou` using org_map
            replace_if_exists(raw_data_df, 'ou', org_map)

            # Replace `dx` using dx_map
            replace_if_exists(raw_data_df, 'dx', dx_map)

            # Replace `co` using option_map
            replace_if_exists(raw_data_df, 'co', option_map)

            # Replace `co_1`, `co_2`, ..., using category_map
            if len(disagg_cols) > 0:
                for col in disagg_cols:
                    replace_if_exists(raw_data_df, col, option_map)

            indicators = raw_data_df["dx"].dropna().unique().tolist()
            periods = sorted(raw_data_df["pe"].dropna().unique().tolist())
            try:
                org_units = (
                    raw_data_df["ou"].dropna().unique().tolist()
                    if "ou" in raw_data_df.columns else raw_data_df["Organisation unit"].dropna().unique().tolist()

                )
            except KeyError:
                org_units = []

            # Rename disaggregation columns to co_1, co_2, etc.
            co_col_map = {}  # e.g. {'co_1': 'cX5k9anHEHd'}
            for idx, col in enumerate(disagg_cols):
                readable_name = category_map.get(col, col)  # Fallback to col if not found
                co_name = f"co_{idx + 1} {readable_name}"
                raw_data_df.rename(columns={col: co_name}, inplace=True)
                co_col_map[co_name] = col

            # Filters

            selected_indicators = st.multiselect("Select indicator(s)", indicators, default=indicators[:1])
            if len(org_units) > 0:
                selected_org_units = st.multiselect("Select org unit(s)", org_units, default=org_units)
            selected_periods = st.multiselect("Select period(s)", periods, default=periods)

            # Multiselect filters for each co_*
            selected_cos = {}
            for co_col in co_col_map:
                co_values = raw_data_df[co_col].dropna().unique().tolist()
                selected_ = st.multiselect(f"Filter by category option: {co_col}", co_values, default=co_values)
                selected_cos[co_col] = selected_

            chart_type = st.selectbox("Select chart type", ["bar", "line", "pie"])

        # Apply all filters
        if len(org_units) > 0:
            # raw_data_df.to_csv("filtered_df_6.csv")
            filtered_df = raw_data_df[
                raw_data_df["dx"].isin(selected_indicators) &
                raw_data_df["pe"].isin(selected_periods)
                ]
            # & raw_data_df[raw_data_df.columns[2]].isin(selected_org_units)

        else:
            # raw_data_df.to_csv("filtered_df_7.csv")
            filtered_df = raw_data_df[
                raw_data_df["dx"].isin(selected_indicators) &
                raw_data_df["pe"].isin(selected_periods)
                ]
        # filtered_df.to_csv("filtered_df_8.csv")
        # Apply co_* filters
        for co_col, selected_vals in selected_cos.items():
            if selected_vals:
                filtered_df = filtered_df[filtered_df[co_col].isin(selected_vals)]

        # filtered_df.to_csv("filtered_df_9.csv")

        if selected_indicators and not filtered_df.empty:
            if chart_backend == "chartjs":
                render_chart_chartjs(filtered_df, selected_indicators, chart_type)
            elif chart_backend == "matplotlib":
                render_chart_matplotlib(filtered_df, selected_indicators, chart_type)
            elif chart_backend == "plotly":
                render_chart_plotly(filtered_df, selected_indicators, chart_type)
        else:
            st.info("No data for the selected filters.")
    else:
        st.warning("No valid data or required columns missing.")

def handle_routing_clarification(router_result: str, state_: dict) -> str:
    """
    Streamlit helper to handle routing decisions.

    Args:
        router_result (str): Output from `routing_decision`.
        state_ (dict): Your agent state, used to persist last_active_agent.

    Returns:
        str: The final agent to route to.
    """
    if router_result != "clarify_agent":
        return router_result  # normal routing

    st.warning("I’m not sure which agent should handle your request. Please choose:")

    # Agent options
    agent_options = {
        "metadata": "Metadata Agent",
        "analytics": "Analytics Agent",
        "data_entry": "Data Entry Agent",
        "event_data": "Event Data Agent",
        "tracker_data": "Tracker Data Agent"
    }

    # Display buttons in Streamlit
    col1_, col2_, col3 = st.columns(3)
    selected_agent = None
    for idx, (agent_id, label_) in enumerate(agent_options.items()):
        col = [col1, col2, col3][idx % 3]
        if col.button(label_):
            selected_agent = agent_id
            state_["last_active_agent"] = agent_id  # persist choice
            st.success(f"Routing to: {label_}")

    # If user hasn’t picked yet, return None
    if selected_agent:
        return selected_agent
    return "clarify_agent"  # still waiting for user input

def handle_suggestion_tool_result(result_):
    suggestion_tool_result_ = result_.get("suggestion_tool_result", "")
    print(f"suggestion_tool_result => {suggestion_tool_result_}")

    if suggestion_tool_result_ and suggestion_tool_result_.get("status") == "multiple_matches":
        suggestions = suggestion_tool_result_["suggestions"]

        # Keep unresolved suggestions in a queue
        if "unresolved_metadata" not in st.session_state:
            st.session_state.unresolved_metadata = []

        st.session_state.unresolved_metadata.append(suggestion_tool_result_)
        st.session_state.show_suggestion_options_data = True
        return True

    return False

# ========== App UI ==========

st.set_page_config(page_title="DHIS2 Assistant", layout="wide")

# Navigation
st.sidebar.header("🧭 Navigation")
page = st.sidebar.selectbox(
    "Choose Interface",
    ["Main Chat", "About"]
)

if page == "Main Chat":
    col1, col2 = st.columns([1, 5])
    with col1:
        st.image("fhi360.png")
    with col2:
        st.title("🗨️ DHIS2 Chat Assistant")

    st.markdown("""
    **Available Agents:**

    1. **📊 Metadata Agent** - Create, update, delete, or retrieve DHIS2 metadata
    2. **📈 Analytics Agent** - Generate charts, reports, trends, and summaries
    3. **📝 Data Entry Agent** - Enter or update aggregate data values
    4. **📅 Event Data Agent** - Handle event-based program data
    5. **👥 Tracker Data Agent** - Manage tracked entity operations

    **Try these examples:**
    - "Create a new dataset for malaria cases"
    - "Show me a trend of ANC visits in Bo district"
    - "Enter 25 malaria cases for Bombali for January"
    - "Process this facility register with Azure Document Intelligence"
    """)

# elif page == "Azure Document Intelligence":
#     st.header("☁️ Azure Document Intelligence")
#     st.markdown("""
#     **Enterprise-grade document processing using Azure Document Intelligence (Form Recognizer).**
#
#     This interface provides:
#     - 🤖 Custom model training for facility register layouts
#     - ✍️ High-accuracy handwriting recognition
#     - 🗺️ Automated DHIS2 mapping and validation
#     - 📊 Production-ready monitoring and analytics
#     """)
#
#     # Check if Azure credentials are configured
#     azure_configured = all([
#         os.getenv("AZURE_FORM_RECOGNIZER_ENDPOINT"),
#         os.getenv("AZURE_FORM_RECOGNIZER_KEY"),
#         os.getenv("AZURE_STORAGE_CONNECTION_STRING")
#     ])
#
#     if azure_configured:
#         st.success("✅ Azure credentials configured")
#
#         # Launch Azure Document Intelligence app
#         st.markdown("""
#         **To access the full Azure Document Intelligence interface, run:**
#         ```bash
#         streamlit run azure_document_intelligence_app.py
#         ```
#         """)
#
#         # Show a preview of features
#         col1, col2, col3 = st.columns(3)
#
#         with col1:
#             st.subheader("📄 Document Processing")
#             st.markdown("""
#             - Upload facility registers
#             - Layout analysis
#             - Custom model processing
#             - Handwriting recognition
#             """)
#
#         with col2:
#             st.subheader("🤖 Model Training")
#             st.markdown("""
#             - Train custom models
#             - Upload training data
#             - Monitor training status
#             - Deploy model versions
#             """)
#
#         with col3:
#             st.subheader("🗺️ DHIS2 Integration")
#             st.markdown("""
#             - Automatic mapping
#             - Data validation
#             - Confidence scoring
#             - DHIS2 submission
#             """)
#     else:
#         st.warning("⚠️ Azure credentials not configured")
#         st.markdown("""
#         **To use Azure Document Intelligence, configure these environment variables:**
#         - `AZURE_FORM_RECOGNIZER_ENDPOINT`
#         - `AZURE_FORM_RECOGNIZER_KEY`
#         - `AZURE_STORAGE_CONNECTION_STRING`
#         """)

elif page == "About":
    st.header("ℹ️ About the System")

    st.markdown("""
    ## 🤖 Multi-Agent Architecture

    This DHIS2 assistant uses a sophisticated multi-agent system with specialized AI agents:

    ### **Available Agents**

    1. **📊 Metadata Agent**
       - Manages DHIS2 metadata (datasets, programs, data elements, org units)
       - Create, update, delete, and retrieve metadata
       - Examples: "Create a new dataset", "What is the UID for ANC visits?"

    2. **📈 Analytics Agent**
       - Generates charts, reports, trends, and summaries
       - Computes totals, averages, and breakdowns
       - Examples: "Show me a trend of malaria cases", "What's the total for Bo district?"

    3. **📝 Data Entry Agent**
       - Handles aggregate data entry and updates
       - Works with standard datasets and data elements
       - Examples: "Enter 25 malaria cases for Bombali for January"

    4. **📅 Event Data Agent**
       - Manages event-based program data
       - Handles single events without registration
       - Examples: "Record a malaria event in Kenema on March 5th"

    5. **👥 Tracker Data Agent**
       - Manages tracked entity operations
       - Handles registrations, follow-ups, and tracked events
       - Examples: "Register a pregnant woman for ANC", "Record a follow-up visit"

    6. **📄 Document Intelligence Agent**
       - Processes scanned documents and facility registers
       - Analyzes document structure and table layouts
       - Maps document columns to DHIS2 metadata
       - Converts handwritten data to DHIS2 JSON payloads

    7. **☁️ Azure Document Intelligence Agent**
       - Enterprise-grade document processing with Azure services
       - Custom model training for facility register layouts
       - High-accuracy handwriting recognition
       - Production-ready document processing pipeline

    ### **How It Works**

    1. **Intelligent Routing**: The system automatically routes your request to the most appropriate agent
    2. **Specialized Processing**: Each agent handles specific types of DHIS2 operations
    3. **Context Awareness**: Agents maintain context and can handle complex multi-step operations
    4. **Error Handling**: Robust error handling and recovery mechanisms

    ### **Technologies Used**

    - **LangGraph**: Multi-agent orchestration and state management
    - **LangChain**: AI agent framework and tool integration
    - **Azure OpenAI**: Large language models for natural language understanding
    - **Azure Document Intelligence**: Enterprise document processing
    - **Streamlit**: User interface and interaction
    - **DHIS2 API**: Integration with DHIS2 systems

    ### **Getting Started**

    1. **Main Chat**: Use the main chat interface for general DHIS2 operations
    2. **Document Intelligence**: Process scanned documents with local OCR
    3. **Azure Document Intelligence**: Enterprise-grade document processing
    4. **About**: Learn more about the system architecture

    ### **Examples**

    **Metadata Operations:**
    - "Create a new dataset for malaria cases"
    - "What data elements are available for ANC visits?"
    - "Update the malaria dataset description"

    **Analytics Queries:**
    - "Show me a trend of malaria cases over the last 6 months"
    - "What's the total number of ANC visits in Bo district?"
    - "Generate a report of immunization coverage"

    **Data Entry:**
    - "Enter 25 malaria cases for Bombali for January 2024"
    - "Update the ANC visits count for Kenema clinic"
    - "Delete the incorrect malaria data for March"

    **Document Processing:**
    - "Process this facility register with Azure Document Intelligence"
    - "Train a custom model for our facility register layout"
    - "Map the extracted data to DHIS2 metadata"
    """)

# Continue with the main chat interface only if "Main Chat" is selected
if page == "Main Chat":
    # ========== Side Bar ==========
    # ===== Initialize session state =====
    if "side_bar_open" not in st.session_state:
        st.session_state.side_bar_open = False
    if "side_bar_open_button" not in st.session_state:
        st.session_state.side_bar_open_button = True

    if st.session_state.get("side_bar_open_button"):
        if st.button("🛠️"):
            st.session_state.side_bar_open = True

    if st.session_state.get("side_bar_open"):
        # Simulate a menu using the sidebar
        st.session_state.side_bar_open_button = False
        with st.sidebar:
            st.header("Settings")
            # if st.button("▶ Run Script"):
            #     # result = run_my_code()
            #     st.success("Button pressed")
            with st.expander("⚙️ Advanced Options"):
                if st.button("Run Metadata Embedder"):
                    with st.spinner("Embedding and indexing metadata..."):
                        run_embedding()
                    st.success("Metadata embedded.")


    upload_file = st.file_uploader(
        "➕ Upload file",
        type=["csv", "xlsx", "xls", "pdf", "png", "jpg", "jpeg", "docx", "doc", "ppt", "pptx"]
    )

    if upload_file is not None:
        file_name = upload_file.name.lower()
        st.write(f"📁 File uploaded: `{file_name}`")

        try:
            if file_name.endswith(".csv"):
                df = pd.read_csv(upload_file)
                st.session_state.raw_data_df_uploaded = df
                st.success("✅ CSV loaded")
                st.dataframe(df)

            elif file_name.endswith((".xlsx", ".xls")):
                df = pd.read_excel(upload_file)
                st.session_state.raw_data_df_uploaded = df
                st.success("✅ Excel loaded")
                st.dataframe(df)

            elif file_name.endswith(".pdf"):
                # Create a named temporary file
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
                    tmp_file.write(upload_file.read())  # write uploaded content to temp file
                    tmp_path = tmp_file.name  # store path to tempfile
                    # Store both the parsed text and metadata in session_state
                    st.session_state.pdf_file = tmp_path
                    st.session_state.pdf_filename = file_name

                    st.success(f"✅ Submitted (original: {file_name})")
                    # st.text_area("📄 PDF Text", full_text, height=300)

            elif file_name.endswith((".png", ".jpg", ".jpeg")):
                image = Image.open(upload_file)
                st.session_state.uploaded_image = image
                st.success("✅ Image loaded")
                st.image(image, caption="Uploaded Image", use_column_width=True)

            elif file_name.endswith(".docx"):
                text = docx2txt.process(upload_file)
                st.session_state.word_text = text
                st.success("✅ DOCX parsed")
                st.text_area("📄 Word Text", text, height=300)

            elif file_name.endswith(".doc"):
                st.warning("Legacy `.doc` support is limited. Please convert to `.docx` for better results.")
                st.session_state.word_text = ""

            elif file_name.endswith(".pptx"):
                prs = Presentation(upload_file)
                slides_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slides_text.append(shape.text)
                full_text = "\n".join(slides_text)
                st.session_state.ppt_text = full_text
                st.success("✅ PowerPoint parsed")
                st.text_area("📄 Slides Text", full_text, height=300)

            elif file_name.endswith(".ppt"):
                st.warning("Legacy `.ppt` support is limited. Please convert to `.pptx` for better results.")
                st.session_state.ppt_text = ""

            else:
                st.warning("Unsupported file type.")

        except Exception as e:
            st.error(f"❌ Error reading file: {e}")

    # ========== Side bar =============

    # # Simulate a menu using the sidebar
    # with st.sidebar:
    #     st.header("Settings")
    #     # if st.button("▶ Run Script"):
    #     #     # result = run_my_code()
    #     #     st.success("Button pressed")
    #     with st.expander("⚙️ Advanced Options"):
    #         if st.button("Run Metadata Embedder"):
    #             with st.spinner("Embedding and indexing metadata..."):
    #                 run_embedding()
    #             st.success("Metadata embedded.")

    # Chat history
    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "suggestions" not in st.session_state:
        st.session_state.suggestions = []
    if "show_chart" not in st.session_state:
        st.session_state.show_chart = False

    if "show_suggestion_options_metadata" not in st.session_state:
        st.session_state.show_suggestion_options_metadata = False

    if "show_suggestion_options_data" not in st.session_state:
        st.session_state.show_suggestion_options_data = False

    for msg in st.session_state.messages:
        role = "assistant" if isinstance(msg, AIMessage) else "user"
        with st.chat_message(role):
            st.markdown(msg.content)

    if prompt := st.chat_input("Ask DHIS2 Assistant..."):
        user_msg = HumanMessage(content=prompt)
        st.session_state.messages.append(user_msg)
        with st.chat_message("user"):
            st.markdown(prompt)

        state = {
            "messages": st.session_state.messages
        }
        raw_df = st.session_state.get("raw_data_df_uploaded", None)
        columns = raw_df.columns.tolist() if raw_df is not None else []
        # Inject column info message FIRST so it's in context before invoking
        messages = st.session_state.get("messages", []).copy()
        if columns:
            column_msg = AIMessage(

                content=f"The uploaded data contains the following columns: {columns}"

            )

            messages.append(column_msg)
            reminder_msg = AIMessage(
                content=(
                    "Reminder: Always call the `submit_aggregate_data` tool for any data submission, update, "
                    "or deletion operations. Do not simulate or fake these actions. "
                    "Confirm success only after the tool is actually called."
                )
            )
            messages.append(reminder_msg)

            state = {

                "messages": messages,
                "dataframe": raw_df,  # Optional: pass the actual dataframe too
                "dataframe_columns": columns,  # Required: LLM needs column names
                "pdf_file": st.session_state.get("pdf_file", None),
                # "image": st.session_state.get("uploaded_image", None),
                # "word_text": st.session_state.get("word_text", None),
                # "ppt_text": st.session_state.get("ppt_text", None)

            }
        route = routing_decision(state)
        # 2. Handle possible clarification
        final_agent = handle_routing_clarification(route, state)

        # 3. Only continue if final_agent is determined
        if final_agent != "clarify_agent":

            if route == "metadata":
                result = metadata_agent_executor.invoke(state)
            elif route == "analytics":
                result = analytics_executor.invoke(state)
                metadata_result = result.get("metadata_result", "")

                if metadata_result and metadata_result.get("status") == "multiple_matches":
                    suggestions = metadata_result["suggestions"]
                    st.session_state.suggestions = suggestions
                    st.session_state.show_suggestion_options_metadata = True


                else:
                    st.session_state.result = result
                    st.session_state.show_chart = True
                    output = result.get("output")

                    assistant_msg = output if isinstance(output, AIMessage) else AIMessage(content=str(output))
                    st.session_state.messages.append(assistant_msg)
                    with st.chat_message("assistant"):
                        st.markdown(assistant_msg.content)

            elif route == "data_entry":
                # Append a strong reminder to the last user message in state before invoking the executor
                # Find the last human message in the conversation history
                # for msg in reversed(state["messages"]):
                #     if isinstance(msg, HumanMessage):
                #         # Append the reminder to the last human message's content
                #         msg.content += (
                #             "\n\nIMPORTANT: You must call the tool `submit_aggregate_data` "
                #             "whenever performing any data submission, update, or deletion. "
                #             "Do not simulate or fake tool calls under any circumstances."
                #         )
                #         break  # Only append once to the most recent human message
                if "raw_data_df_uploaded" not in st.session_state:
                    st.session_state.file_uploaded = False
                else:
                    st.session_state.file_uploaded = True
                result = data_entry_executor.invoke(state)

                suggestion_tool_result = result.get("suggestion_tool_result", "")
                print(f"suggestion_tool_result 1 => {suggestion_tool_result}")

                if suggestion_tool_result and suggestion_tool_result.get("status") == "multiple_matches":
                    suggestions = suggestion_tool_result["suggestions"]
                    st.session_state.suggestions = suggestions
                    st.session_state.show_suggestion_options_data = True


                else:
                    st.session_state.result = result
                    output = result.get("output")
                    assistant_msg = output if isinstance(output, AIMessage) else AIMessage(content=str(output))
                    st.session_state.messages.append(assistant_msg)
                    with st.chat_message("assistant"):
                        st.markdown(assistant_msg.content)


            elif route == "event_data":
                result = event_data_executor.invoke(state)
            elif route == "tracker_data":
                result = tracker_data_executor.invoke(state)
                # print(result)
            # elif route == "azure_document_intelligence":
            #     result = azure_document_intelligence_executor.invoke(state)
            else:
                error_msg = AIMessage(content="❌ Unknown routing decision.")
                st.session_state.messages.append(error_msg)
                st.chat_message("assistant").markdown(error_msg.content)
                st.stop()

            if len(st.session_state.suggestions) < 1:
                output = result.get("output")
                assistant_msg = output if isinstance(output, AIMessage) else AIMessage(content=str(output))
                st.session_state.messages.append(assistant_msg)
                with st.chat_message("assistant"):
                    st.markdown(assistant_msg.content)

    if st.session_state.get("show_chart", False):
        # print(f"The show_chart {st.session_state.get("show_chart")}")
        st.info(f"The show_chart {st.session_state.get("show_chart")}")
        chart_data(st.session_state.result, "chartjs")

    # Show metadata options after user prompt
    if st.session_state.get("show_suggestion_options_metadata", False):
        with st.chat_message("assistant"):
            st.markdown("There are multiple metadata matches. Please choose one:")

            for i, option in enumerate(st.session_state.suggestions):
                label = f"{option['name']} ({option['doc_type']})"
                if st.button(label, key=f"metadata_option_{i}"):
                    st.session_state.selected_metadata = option  # Store full dict (id, name, doc_type, score)
                    st.session_state.trigger_metadata_retry = True  # 🔁 trigger new processing
                    st.session_state.show_suggestion_options_metadata = False
                    st.rerun()

    # Show metadata options after user prompt
    if st.session_state.get("show_suggestion_options_data", False):
        with st.chat_message("assistant"):
            st.markdown("There are multiple metadata matches for data processing. Please choose one:")

            for i, option in enumerate(st.session_state.suggestions):
                label = f"{option['name']} ({option['doc_type']})"
                if st.button(label, key=f"data_entry_option_{i}"):
                    st.session_state.selected_metadata = option  # Store full dict (id, name, doc_type, score)
                    st.session_state.trigger_data_retry = True  # 🔁 trigger new processing
                    st.session_state.show_suggestion_options_data = False
                    st.rerun()


    if st.session_state.get("trigger_data_retry"):
        selected = st.session_state.selected_metadata

        # Reconstruct the last actual user question
        original_user_msg = ""
        for msg in reversed(st.session_state.messages):
            if isinstance(msg, HumanMessage):
                original_user_msg = msg.content
                break

        # Build a clearer retry message with metadata reference
        augmented_msg = HumanMessage(
            # content=f"{original_user_msg} (selected: {selected['id']} - {selected['name']})"
            content=f"{selected['doc_type']} selected:  {selected['name']} -  {selected['id']})"

        )

        st.session_state.messages.append(augmented_msg)

        with st.chat_message("user"):
            st.markdown(original_user_msg)  # ✅ Show only original question in UI

        # Re-invoke agent
        state = {"messages": st.session_state.messages}
        result = data_entry_executor.invoke(state)

        suggestion_tool_result = result.get("suggestion_tool_result", "")
        print(f"suggestion_tool_result 1 => {suggestion_tool_result}")

        if suggestion_tool_result and suggestion_tool_result.get("status") == "multiple_matches":
            suggestions = suggestion_tool_result["suggestions"]
            st.session_state.suggestions = suggestions
            st.session_state.show_suggestion_options_data = True


        else:
            st.session_state.result = result
            output = result.get("output")
            assistant_msg = output if isinstance(output, AIMessage) else AIMessage(content=str(output))
            st.session_state.messages.append(assistant_msg)
            with st.chat_message("assistant"):
                st.markdown(assistant_msg.content)






        st.session_state.result = result
        st.session_state.trigger_data_retry = False

        # output = result.get("output")
        # assistant_msg = output if isinstance(output, AIMessage) else AIMessage(content=str(output))
        # st.session_state.messages.append(assistant_msg)
        #
        # with st.chat_message("assistant"):
        #     st.markdown(assistant_msg.content)

        st.rerun()  # ✅ Force chart to render

    if st.session_state.get("trigger_metadata_retry"):
        selected = st.session_state.selected_metadata

        # Reconstruct the last actual user question
        original_user_msg = ""
        for msg in reversed(st.session_state.messages):
            if isinstance(msg, HumanMessage):
                original_user_msg = msg.content
                break

        # Build a clearer retry message with metadata reference
        augmented_msg = HumanMessage(
            # content=f"{original_user_msg} (selected: {selected['id']} - {selected['name']})"
            content=f"{selected['doc_type']} selected:  {selected['name']} -  {selected['id']})"

        )

        st.session_state.messages.append(augmented_msg)

        with st.chat_message("user"):
            st.markdown(original_user_msg)  # ✅ Show only original question in UI

        # Re-invoke agent
        state = {"messages": st.session_state.messages}
        result = analytics_executor.invoke(state)

        st.session_state.result = result
        st.session_state.show_chart = True
        st.session_state.trigger_metadata_retry = False

        output = result.get("output")
        assistant_msg = output if isinstance(output, AIMessage) else AIMessage(content=str(output))
        st.session_state.messages.append(assistant_msg)

        with st.chat_message("assistant"):
            st.markdown(assistant_msg.content)

        st.rerun()  # ✅ Force chart to render



